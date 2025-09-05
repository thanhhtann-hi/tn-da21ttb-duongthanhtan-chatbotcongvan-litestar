# file: src/modules/chat/service/text_extract.py
# updated: 2025-08-31 (v1.3.1)
# purpose:
#   - Trích xuất văn bản thô từ: DOC/DOCX/RTF/ODT, XLS/XLSX/CSV/ODS,
#     PPT/PPTX/ODP, TXT/SVG và nhiều file mã nguồn phổ biến (kể cả .ipynb).
#   - Phụ thuộc ngoài đều là "tùy chọn" (try/except). Không có lib → fallback an toàn.
#   - API tối giản: extract_text(path, ext=None) -> SimpleResult(ok, text, pages_text, total_pages, meta)
#
# lưu ý:
#   - "pages_text" có ý nghĩa với loại nhiều trang (PPTX/ODP: mỗi slide; XLS/XLSX: mỗi sheet).
#   - Có giới hạn đọc kích thước để tránh ngốn RAM.
#   - Không xử lý PDF ở đây (PDF thuộc OCR module riêng).

from __future__ import annotations

import os
import re
import csv
import json
import zipfile
from dataclasses import dataclass
from typing import Optional, List, Dict, Iterable

__all__ = [
    "extract_text", "SimpleResult",
    "DOCS", "SHEETS", "SLIDES", "CODE", "OTHERS",
]

# ──────────────────────────────────────────────────────────────────────────────
# Cấu hình (có thể tinh chỉnh bằng ENV)
# ──────────────────────────────────────────────────────────────────────────────
def _env_int(name: str, default: int) -> int:
    try:
        return int(os.getenv(name, str(default)).strip())
    except Exception:
        return default

# Giới hạn chung
READ_MAX_BYTES           = _env_int("TEXTEX_MAX_TEXT_BYTES", 16 * 1024 * 1024)  # 16 MB
CSV_ROW_LIMIT            = _env_int("TEXTEX_CSV_ROW_LIMIT", 50000)
XLS_CELL_LIMIT           = _env_int("TEXTEX_XLS_CELL_LIMIT", 20000)
XLSX_CELL_LIMIT          = _env_int("TEXTEX_XLSX_CELL_LIMIT", 20000)
PPTX_SLIDE_LIMIT         = _env_int("TEXTEX_PPTX_SLIDE_LIMIT", 500)
CSV_SNIFF_BYTES          = _env_int("TEXTEX_CSV_SNIFF_BYTES", 128 * 1024)
ODF_XML_STRIP            = True  # gọn XML → text

@dataclass
class SimpleResult:
    ok: bool
    text: str
    pages_text: Optional[List[str]] = None
    total_pages: Optional[int] = None
    meta: Optional[Dict[str, str]] = None


# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────
def _is_probably_binary(chunk: bytes) -> bool:
    if not chunk:
        return False
    # nếu có nhiều bytes điều khiển → coi là nhị phân
    ctrl = sum(1 for b in chunk if (b < 9 or 13 < b < 32))
    return (ctrl / max(1, len(chunk))) > 0.3

def _read_text_file(path: str, limit_bytes: int = READ_MAX_BYTES) -> str:
    """Đọc file text/ mã nguồn an toàn (UTF-8, ignore-errors), có giới hạn."""
    size = 0
    bufs: List[str] = []
    with open(path, "rb") as f:
        while True:
            if size >= limit_bytes:
                break
            chunk = f.read(min(1 << 18, limit_bytes - size))  # 256KB/đợt
            if not chunk:
                break
            size += len(chunk)
            # decode bất kể có nhị phân rải rác
            bufs.append(chunk.decode("utf-8", errors="ignore"))
    return "".join(bufs)

def _strip_xml(txt: str) -> str:
    # đơn giản: bỏ tag → thu gọn whitespace
    txt = re.sub(r"<[^>]+>", " ", txt)
    txt = re.sub(r"\s+", " ", txt)
    return txt.strip()

def _unzip_read(path: str, inner: str, limit_bytes: int = READ_MAX_BYTES) -> str:
    """Đọc một entry text từ file ZIP (OOXML/ODF). Giới hạn kích thước để tránh RAM spike."""
    try:
        with zipfile.ZipFile(path) as z:
            if inner not in z.namelist():
                return ""
            with z.open(inner) as f:
                size = 0
                chunks: List[bytes] = []
                while True:
                    if size >= limit_bytes:
                        break
                    b = f.read(min(1 << 18, limit_bytes - size))
                    if not b:
                        break
                    chunks.append(b)
                    size += len(b)
                return b"".join(chunks).decode("utf-8", errors="ignore")
    except Exception:
        return ""

def _join_nonempty(parts: Iterable[str], sep: str = "\n") -> str:
    return sep.join([p for p in parts if p and p.strip()])

def _normalize_spaces(s: str) -> str:
    return re.sub(r"[ \t]+\n", "\n", re.sub(r"[ \t]{2,}", " ", s or "")).strip()

def _mk_result(ok: bool, text: str, *, pages: Optional[List[str]] = None,
               kind: str = "", extra_meta: Optional[Dict[str, str]] = None) -> SimpleResult:
    meta = {"kind": kind} if kind else {}
    if extra_meta:
        meta.update({k: str(v) for k, v in extra_meta.items()})
    total = len(pages) if pages is not None else None
    return SimpleResult(ok=ok, text=text, pages_text=pages, total_pages=total, meta=meta)


# ──────────────────────────────────────────────────────────────────────────────
# DOCX / DOC / RTF / ODT
# ──────────────────────────────────────────────────────────────────────────────
def _docx(path: str) -> str:
    # Ưu tiên python-docx → fallback docx2txt → fallback unzip
    # paragraphs + (optional) tables text
    try:
        import docx  # python-docx
        d = docx.Document(path)
        blocks: List[str] = []

        # paragraphs
        blocks.extend([p.text for p in d.paragraphs if p.text and p.text.strip()])

        # tables (đơn giản)
        for t in getattr(d, "tables", []):
            for row in t.rows:
                cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                if cells:
                    blocks.append(" | ".join(cells))

        return _normalize_spaces(_join_nonempty(blocks, "\n"))
    except Exception:
        pass

    try:
        import docx2txt  # type: ignore
        txt = docx2txt.process(path) or ""
        return _normalize_spaces(txt)
    except Exception:
        pass

    # unzip: word/document.xml + headers/footers (thường chứa nội dung)
    xmls = []
    for candidate in [
        "word/document.xml",
        "word/header1.xml", "word/header2.xml", "word/header3.xml",
        "word/footer1.xml", "word/footer2.xml", "word/footer3.xml",
    ]:
        xml = _unzip_read(path, candidate)
        if xml:
            xmls.append(xml)
    return _normalize_spaces(_strip_xml("\n".join(xmls))) if xmls else ""


def _doc(path: str) -> str:
    # DOC binary: phụ thuộc ngoài (textract / antiword). Nếu không có → trả rỗng.
    try:
        import textract  # type: ignore
        raw = textract.process(path)  # có thể hơi nặng → người dùng tự cài
        return _normalize_spaces((raw.decode("utf-8", errors="ignore") or ""))
    except Exception:
        # cố nốt antiword nếu có
        try:
            import subprocess, shlex  # noqa: F401
            out = subprocess.check_output(["antiword", path], stderr=subprocess.DEVNULL)
            return _normalize_spaces(out.decode("utf-8", errors="ignore"))
        except Exception:
            return ""


def _rtf(path: str) -> str:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
        raw = _read_text_file(path, limit_bytes=min(READ_MAX_BYTES, 32 * 1024 * 1024))
        return _normalize_spaces(rtf_to_text(raw) or "")
    except Exception:
        # fallback thô: bỏ điều khiển rtf và dấu {}
        raw = _read_text_file(path, limit_bytes=min(READ_MAX_BYTES, 32 * 1024 * 1024))
        raw = re.sub(r"\\[a-z]+\d*", " ", raw)
        raw = re.sub(r"[{}]", " ", raw)
        return _normalize_spaces(raw)


def _odt(path: str) -> str:
    xml = _unzip_read(path, "content.xml")
    return _normalize_spaces(_strip_xml(xml)) if xml and ODF_XML_STRIP else _normalize_spaces(xml)


# ──────────────────────────────────────────────────────────────────────────────
# XLS/XLSX/CSV/ODS
# ──────────────────────────────────────────────────────────────────────────────
def _xlsx(path: str, cell_limit: int = XLSX_CELL_LIMIT) -> SimpleResult:
    """openpyxl: sheet → dòng ' | ' cell; pages_text = mỗi sheet 1 trang."""
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        pages: List[str] = []
        for ws in wb.worksheets:
            lines: List[str] = [f"[Sheet] {ws.title}"]
            count = 0
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) for v in row if v is not None and str(v).strip() != ""]
                if not vals:
                    continue
                line = " | ".join(vals)
                lines.append(line)
                count += len(vals)
                if count >= cell_limit:
                    lines.append("…")
                    break
            page_text = _normalize_spaces("\n".join(lines))
            pages.append(page_text)
        full = _normalize_spaces("\n".join(pages))
        return _mk_result(bool(full), full, pages=pages, kind="xlsx", extra_meta={"sheets": str(len(pages))})
    except Exception:
        return _mk_result(False, "", kind="xlsx")


def _xls(path: str, cell_limit: int = XLS_CELL_LIMIT) -> SimpleResult:
    """xlrd: sheet → dòng; pages_text = mỗi sheet 1 trang."""
    try:
        import xlrd  # type: ignore
        wb = xlrd.open_workbook(path)
        pages: List[str] = []
        for i in range(wb.nsheets):
            sh = wb.sheet_by_index(i)
            lines: List[str] = [f"[Sheet] {sh.name}"]
            count = 0
            for r in range(sh.nrows):
                vals = [str(v) for v in sh.row_values(r) if v not in ("", None)]
                if not vals:
                    continue
                line = " | ".join(vals)
                lines.append(line)
                count += len(vals)
                if count >= cell_limit:
                    lines.append("…")
                    break
            page_text = _normalize_spaces("\n".join(lines))
            pages.append(page_text)
        full = _normalize_spaces("\n".join(pages))
        return _mk_result(bool(full), full, pages=pages, kind="xls", extra_meta={"sheets": str(len(pages))})
    except Exception:
        return _mk_result(False, "", kind="xls")


def _sniff_csv_dialect(sample: str) -> csv.Dialect | None:
    try:
        return csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except Exception:
        return None


def _csv(path: str, row_limit: int = CSV_ROW_LIMIT) -> str:
    out: List[str] = []
    # thử sniff delimiter
    dialect = None
    try:
        with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
            head = f.read(CSV_SNIFF_BYTES)
            dialect = _sniff_csv_dialect(head)
    except Exception:
        dialect = None

    try:
        with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.reader(f, dialect=dialect) if dialect else csv.reader(f)
            for i, row in enumerate(reader):
                out.append(" | ".join([c.strip() for c in row]))
                if i + 1 >= row_limit:
                    out.append("…")
                    break
    except Exception:
        return ""
    return _normalize_spaces("\n".join(out))


def _ods(path: str) -> str:
    xml = _unzip_read(path, "content.xml")
    return _normalize_spaces(_strip_xml(xml)) if xml and ODF_XML_STRIP else _normalize_spaces(xml)


# ──────────────────────────────────────────────────────────────────────────────
# PPT/PPTX/ODP
# ──────────────────────────────────────────────────────────────────────────────
def _pptx(path: str, slide_limit: int = PPTX_SLIDE_LIMIT) -> SimpleResult:
    """python-pptx: pages_text = mỗi slide là một phần tử."""
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(path)
        pages: List[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            if i > slide_limit:
                pages.append("…")
                break
            parts: List[str] = [f"[Slide {i}]"]
            for shp in slide.shapes:
                try:
                    if hasattr(shp, "text") and shp.text:
                        parts.append(shp.text.strip())
                except Exception:
                    continue
            # ghi chú speaker notes (nếu có)
            try:
                if slide.has_notes_slide and slide.notes_slide:
                    note_frame = slide.notes_slide.notes_text_frame
                    if note_frame and note_frame.text:
                        parts.append(f"[Notes] {note_frame.text.strip()}")
            except Exception:
                pass
            page_text = _normalize_spaces("\n".join(parts))
            pages.append(page_text)
        full = _normalize_spaces("\n\n".join(pages))
        return _mk_result(bool(full), full, pages=pages, kind="pptx", extra_meta={"slides": str(len(pages))})
    except Exception:
        return _mk_result(False, "", kind="pptx")


def _ppt(path: str) -> str:
    # Không có thư viện ổn định & nhẹ cho PPT binary → trả rỗng (tránh mang thêm phụ thuộc nặng).
    return ""


def _odp(path: str) -> SimpleResult:
    """
    ODP (OpenDocument Presentation):
    - Tách sơ bộ theo <draw:page ...> để tạo pages_text.
    - Nếu không tách được: strip XML toàn bộ.
    """
    xml = _unzip_read(path, "content.xml")
    if not xml:
        return _mk_result(False, "", kind="odp")
    try:
        # Chia theo thẻ slide
        chunks = re.split(r"<draw:page\b[^>]*>", xml, flags=re.IGNORECASE)
        pages: List[str] = []
        for ch in chunks:
            if not ch:
                continue
            # cắt tại kết thúc page (heuristic)
            ch2 = ch.split("</draw:page>", 1)[0]
            txt = _strip_xml(ch2).strip()
            if txt:
                pages.append(_normalize_spaces(txt))
        # thêm nhãn slide cho đồng nhất
        pages = [f"[Slide {i+1}] {p}" for i, p in enumerate(pages)]
        if pages:
            full = _normalize_spaces("\n\n".join(pages))
            return _mk_result(True, full, pages=pages, kind="odp", extra_meta={"slides": str(len(pages))})
        else:
            # không tách được slide → strip toàn bộ
            return _mk_result(True, _normalize_spaces(_strip_xml(xml)), pages=None, kind="odp")
    except Exception:
        return _mk_result(True, _normalize_spaces(_strip_xml(xml)), pages=None, kind="odp")


# ──────────────────────────────────────────────────────────────────────────────
# SVG (lấy text trong <text>… hoặc <tspan>…)
# ──────────────────────────────────────────────────────────────────────────────
def _svg(path: str) -> str:
    try:
        raw = _read_text_file(path, limit_bytes=min(READ_MAX_BYTES, 8 * 1024 * 1024))
        # ưu tiên nội dung trong thẻ <text>/<tspan>
        texts = re.findall(r"<(?:text|tspan)[^>]*>(.*?)</(?:text|tspan)>", raw, flags=re.IGNORECASE | re.DOTALL)
        if texts:
            body = " ".join(_strip_xml(t) for t in texts)
        else:
            body = _strip_xml(raw)
        return _normalize_spaces(body)
    except Exception:
        return ""


# ──────────────────────────────────────────────────────────────────────────────
# IPYNB (Notebook) — trích Markdown + code cells (gọn nhẹ, không phụ thuộc ngoài)
# ──────────────────────────────────────────────────────────────────────────────
def _ipynb(path: str) -> str:
    try:
        raw = _read_text_file(path, limit_bytes=READ_MAX_BYTES)
        nb = json.loads(raw)
        out: List[str] = []
        for cell in nb.get("cells", []):
            ctype = cell.get("cell_type")
            if ctype == "markdown":
                src = cell.get("source", [])
                out.append(_normalize_spaces("".join(src)))
            elif ctype == "code":
                src = cell.get("source", [])
                # giữ code, nối dòng
                code = "".join(src)
                if code.strip():
                    out.append(code.strip())
        return _normalize_spaces("\n\n".join(out))
    except Exception:
        return ""


# ──────────────────────────────────────────────────────────────────────────────
# Public API
# ──────────────────────────────────────────────────────────────────────────────
DOCS   = {"doc", "docx", "rtf", "odt"}
SHEETS = {"xls", "xlsx", "csv", "ods"}
SLIDES = {"ppt", "pptx", "odp"}
CODE = {
    "js","ts","jsx","tsx","json","py","rb","php","java","kt","kts","c","h","cpp","hpp","cs","go","rs",
    "html","htm","xml","css","scss","sass","less","md","markdown","sql","sh","bash","zsh","bat","ps1",
    "yml","yaml","toml","ini","gradle","m","mm","swift","dart","lua","r","pl","ipynb"
}
OTHERS = {"txt", "svg", "log"}

def extract_text(path: str, ext: Optional[str] = None) -> SimpleResult:
    """
    Trả về SimpleResult:
        - ok: bool
        - text: str (toàn bộ text rút ra — có thể rỗng)
        - pages_text: List[str] hoặc None (nếu có khái niệm 'trang')
        - total_pages: int hoặc None
        - meta: dict phụ (kind/…)
    Không raise exception ra ngoài.
    """
    e = (ext or os.path.splitext(path)[1][1:] or "").lower()

    try:
        if e in {"docx"}:
            txt = _docx(path)
            return _mk_result(bool(txt), txt, kind="docx")
        elif e in {"doc"}:
            txt = _doc(path)
            return _mk_result(bool(txt), txt, kind="doc")
        elif e in {"rtf"}:
            txt = _rtf(path)
            return _mk_result(bool(txt), txt, kind="rtf")
        elif e in {"odt"}:
            txt = _odt(path)
            return _mk_result(bool(txt), txt, kind="odt")

        elif e in {"xlsx"}:
            return _xlsx(path)
        elif e in {"xls"}:
            return _xls(path)
        elif e in {"csv"}:
            txt = _csv(path)
            return _mk_result(bool(txt), txt, kind="csv")
        elif e in {"ods"}:
            txt = _ods(path)
            return _mk_result(bool(txt), txt, kind="ods")

        elif e in {"pptx"}:
            return _pptx(path)
        elif e in {"ppt"}:
            txt = _ppt(path)
            return _mk_result(bool(txt), txt, kind="ppt")
        elif e in {"odp"}:
            return _odp(path)

        elif e in {"svg"}:
            txt = _svg(path)
            return _mk_result(bool(txt), txt, kind="svg")

        elif e in {"ipynb"}:
            txt = _ipynb(path)
            return _mk_result(bool(txt), txt, kind="ipynb")

        elif e in CODE or e in {"txt", "log"}:
            txt = _read_text_file(path, limit_bytes=READ_MAX_BYTES)
            return _mk_result(bool(txt), txt, kind="text")

        # không nhận diện được → nếu là text thuần thì vẫn đọc
        try:
            txt = _read_text_file(path, limit_bytes=READ_MAX_BYTES)
            return _mk_result(bool(txt), txt, kind="text-unknown")
        except Exception:
            return _mk_result(False, "", kind="unknown")

    except Exception:
        return _mk_result(False, "", kind=(e or "unknown"))
